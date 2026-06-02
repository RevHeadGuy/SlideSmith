"""
PDF Extraction Node - Extracts content from PDF and generates outline
"""

import json
import re
import ollama
from pypdf import PdfReader
from typing import Optional
from outline_node import PresentationState


# ── Shared robust JSON parser (same logic as outline_node + ppt_node) ────────

def _clean_json(text: str) -> str:
    """Strip trailing commas before ] or } so LLM output parses cleanly."""
    return re.sub(r',\s*([\]\}])', r'\1', text)


def _extract_json_object(text: str) -> str:
    """Pull out the first complete {...} block — handles LLM preamble text."""
    start = text.find('{')
    if start == -1:
        return text
    depth = 0
    for i in range(start, len(text)):
        if text[i] == '{':
            depth += 1
        elif text[i] == '}':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]
    return text[start:]


def _try_parse(text: str):
    """Try json.loads with progressively more aggressive cleanup."""
    # Strip code fences first
    raw = text.strip()
    if '```' in raw:
        s = raw.find('```json') + 7 if '```json' in raw else raw.find('```') + 3
        e = raw.rfind('```')
        if e > s:
            raw = raw[s:e].strip()

    for candidate in (raw, _clean_json(raw)):
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass

    # Extract JSON object block and retry
    extracted = _extract_json_object(raw)
    for candidate in (extracted, _clean_json(extracted)):
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass

    return None


# ── Shared outline-from-text helper ──────────────────────────────────────────

def _generate_outline_from_text(
    state: PresentationState,
    source_text: str,
    source_label: str,          # "PDF" or "PowerPoint"
    error_status: str,          # "pdf_error" or "pptx_error"
) -> PresentationState:
    """
    Call Llama3 to turn extracted text into a structured JSON outline.
    Shared by both PDF and PPTX paths.
    """
    try:
        max_chars = 3000
        if len(source_text) > max_chars:
            source_text = source_text[:max_chars] + "\n... (content truncated)"

        prompt = f"""Based on the following {source_label} content, create a professional {state['slide_count']}-slide presentation outline.

{source_label} Content:
{source_text}

IMPORTANT: Return ONLY valid JSON — no markdown, no code blocks, no explanation.

{{
    "title": "Presentation Title (derived from content)",
    "slides": [
        {{
            "slide_number": 1,
            "title": "Slide Title",
            "description": "2-3 sentence paragraph about this slide.",
            "key_points": ["point 1", "point 2", "point 3"]
        }}
    ],
    "summary": "One sentence summary"
}}

Include exactly {state['slide_count']} slides. Each slide must have slide_number, title, description, and key_points."""

        print(f"DEBUG: Generating outline from {source_label} content")

        response = ollama.generate(
            model="llama3",
            prompt=prompt,
            stream=False,
            options={"num_predict": 2048},
        )

        response_text = response['response']
        print(f"DEBUG: Ollama response received (length: {len(response_text)})")

        outline = _try_parse(response_text)

        if outline and 'slides' in outline and isinstance(outline['slides'], list):
            state['content'] = outline
            state['status'] = 'outline_generated'
            print(f"DEBUG: Successfully parsed JSON outline from {source_label} ({len(outline['slides'])} slides)")
        else:
            # Last resort — store raw so ppt_node can attempt extraction
            print(f"DEBUG: Could not parse JSON from {source_label} response, storing raw")
            state['content'] = {
                "title": f"{source_label}-based Presentation",
                "outline": response_text,
                "status": "raw",
            }
            state['status'] = 'outline_generated'

        state['generated_at'] = __import__('datetime').datetime.now().isoformat()
        return state

    except Exception as e:
        print(f"ERROR in _generate_outline_from_text ({source_label}): {str(e)}")
        import traceback
        traceback.print_exc()
        state['status'] = error_status
        state['error'] = f"Failed to generate outline from {source_label}: {str(e)}"
        return state


# ── Public node functions ─────────────────────────────────────────────────────

def extract_pdf_text(pdf_path: str) -> str:
    try:
        reader = PdfReader(pdf_path)
        text = ""
        page_count = len(reader.pages)
        print(f"DEBUG: PDF has {page_count} pages")
        for page_num, page in enumerate(reader.pages):
            text += f"\n--- Page {page_num + 1} ---\n"
            text += page.extract_text() or ""
        print(f"DEBUG: Extracted {len(text)} characters from PDF")
        return text
    except Exception as e:
        print(f"ERROR: Failed to extract PDF text: {str(e)}")
        raise


def extract_pdf_node(state: PresentationState) -> PresentationState:
    try:
        if not state.get('pdf_path'):
            raise ValueError("No PDF path provided")
        pdf_path = state['pdf_path']
        print(f"DEBUG: Extracting text from PDF: {pdf_path}")
        pdf_text = extract_pdf_text(pdf_path)
        state['pdf_content'] = pdf_text
        state['status'] = 'pdf_extracted'
        if not state.get('topic') or state['topic'].strip() == '':
            state['topic'] = "PDF Document"
        return state
    except Exception as e:
        print(f"ERROR in extract_pdf_node: {str(e)}")
        import traceback
        traceback.print_exc()
        state['status'] = 'pdf_error'
        state['error'] = f"Failed to extract PDF: {str(e)}"
        return state


def generate_outline_from_pdf(state: PresentationState) -> PresentationState:
    if not state.get('pdf_content'):
        state['status'] = 'pdf_error'
        state['error'] = "No PDF content available for outline generation"
        return state
    return _generate_outline_from_text(
        state,
        source_text=state['pdf_content'],
        source_label="PDF",
        error_status="pdf_error",
    )


def extract_pptx_text(pptx_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        print(f"DEBUG: Extracted {len(text)} characters from PPTX")
        return text
    except Exception as e:
        print(f"ERROR: Failed to extract PPTX text: {str(e)}")
        raise


def extract_pptx_node(state: PresentationState) -> PresentationState:
    try:
        if not state.get('pdf_path'):
            raise ValueError("No PPTX path provided")
        pptx_path = state['pdf_path']
        print(f"DEBUG: Extracting text from PPTX: {pptx_path}")
        pptx_text = extract_pptx_text(pptx_path)
        state['pdf_content'] = pptx_text
        state['status'] = 'pdf_extracted'
        if not state.get('topic') or state['topic'].strip() == '':
            state['topic'] = "PowerPoint Presentation"
        return state
    except Exception as e:
        print(f"ERROR in extract_pptx_node: {str(e)}")
        import traceback
        traceback.print_exc()
        state['status'] = 'pptx_error'
        state['error'] = f"Failed to extract PPTX: {str(e)}"
        return state


def generate_outline_from_pptx(state: PresentationState) -> PresentationState:
    if not state.get('pdf_content'):
        state['status'] = 'pptx_error'
        state['error'] = "No PPTX content available for outline generation"
        return state
    return _generate_outline_from_text(
        state,
        source_text=state['pdf_content'],
        source_label="PowerPoint",
        error_status="pptx_error",
    )
