"""
Presentation Generator - Creates PPTX files aligned with SlideSmith theme
Uses: python-pptx, Pillow
Supports themes: modern, dark, minimal
"""

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches
from pathlib import Path
from typing import List, Dict, Tuple

import re as _re

def format_large_number(num: float) -> str:
    """
    Format large numbers with abbreviations (M, B, T, K).
    Examples: 2500000 -> "2.5M", 1200000000 -> "1.2B"
    """
    if num >= 1_000_000_000:
        return f"{num / 1_000_000_000:.1f}B".rstrip('0').rstrip('.')
    elif num >= 1_000_000:
        return f"{num / 1_000_000:.1f}M".rstrip('0').rstrip('.')
    elif num >= 1_000:
        return f"{num / 1_000:.1f}K".rstrip('0').rstrip('.')
    else:
        return f"{num:.0f}"

def extract_time_series_data(key_points: List[str], description: str = "") -> Tuple[List[str], List[float]]:
    """
    Extract time series data (trends) from key_points.
    Pattern: "Year/Quarter/Month: value" or "value in Year/Quarter"
    Returns (labels, values) tuples
    """
    results = []
    
    # Pattern: "2020: 100", "Q1 2023: 250", "Jan 2024: 45"
    time_pattern = _re.compile(
        r'((?:20|21)\d{2}|Q[1-4]|Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[\s:]?(?:20|21)?(\d{2})?[\s:\-–]*(\d+(?:\.\d+)?)\s*(million|billion|thousand|k|m|b)?',
        _re.IGNORECASE
    )
    
    for pt in key_points[:10]:
        for match in time_pattern.finditer(pt):
            label = match.group(1)
            value = float(match.group(3))
            multiplier = match.group(4)
            
            if multiplier:
                mult = multiplier.lower()
                if mult in ('billion', 'b'): value *= 1_000_000_000
                elif mult in ('million', 'm'): value *= 1_000_000
                elif mult in ('thousand', 'k'): value *= 1_000
            
            results.append((label, value))
    
    if results:
        labels = [r[0] for r in results]
        values = [r[1] for r in results]
        return labels, values
    return [], []

def extract_comparison_series(key_points: List[str], description: str = "") -> Dict[str, List[Tuple[str, float]]]:
    """
    Extract multiple series for comparison charts.
    Pattern: "Category A: Q1=100 Q2=150 Q3=200" or similar structures
    Returns dict like {"Series 1": [("Q1", 100), ("Q2", 150)], ...}
    """
    series_dict = {}
    all_text = " ".join(key_points) + " " + description
    
    # Look for patterns like "Product A: Jan 50, Feb 75, Mar 120"
    series_pattern = _re.compile(
        r'([A-Za-z][A-Za-z0-9\s&/\-]{2,20}?)[\s:\-]+((?:[A-Za-z]{3}\s+\d+[,;]?\s*)+)',
        _re.IGNORECASE
    )
    
    for match in series_pattern.finditer(all_text):
        series_name = match.group(1).strip()
        values_text = match.group(2)
        
        # Parse individual values
        value_pattern = _re.compile(r'([A-Za-z]{3})\s+(\d+(?:\.\d+)?)')
        values = []
        for val_match in value_pattern.finditer(values_text):
            label = val_match.group(1)
            value = float(val_match.group(2))
            values.append((label, value))
        
        if len(values) >= 2:
            series_dict[series_name] = values
    
    return series_dict if series_dict else {}

def has_numeric_content(slide: Dict) -> bool:
    """Return True if slide key_points contain numbers/percentages/years."""
    all_text = ' '.join(slide.get('key_points', [])) + ' ' + slide.get('description', '')
    return bool(_re.search(r'\b(\d+%|\$[\d,]+|\d{4}|\d+\s*(million|billion|thousand|x|times)|\b\d{2,}\b)', all_text, _re.IGNORECASE))

def extract_pie_data(key_points: List[str], description: str = ""):
    """
    Extract (label, value) pairs from key_points AND description that contain percentages.
    Returns list of (label, float) or empty list if fewer than 2 found.
    e.g. "iPhone segment, accounting for 52% of total revenue" -> ("iPhone segment", 52.0)
    """
    results = []
    seen_values = set()

    # Pattern 1: "Label: 45%" or "Label - 45%"
    pct_labeled = _re.compile(r'([\w\s&/]+?)[\s:\-–]+(\d+(?:\.\d+)?)\s*%', _re.IGNORECASE)
    # Pattern 2: "52% of total revenue" or "52% revenue"
    pct_inline  = _re.compile(r'(\d+(?:\.\d+)?)\s*%\s+(?:of\s+)?([a-zA-Z][\w\s]{2,30}?)(?:[,\.]|$)', _re.IGNORECASE)

    def _add(label, value):
        label = label.strip().strip('•:–-').strip()
        # Deduplicate by value
        if value not in seen_values and label and 0 < value <= 100:
            seen_values.add(value)
            results.append((label[:40], value))

    # Scan key_points first
    for pt in key_points:
        m = pct_labeled.match(pt.strip().strip('•').strip())
        if m:
            _add(m.group(1), float(m.group(2)))
            continue
        m2 = pct_inline.search(pt)
        if m2:
            _add(m2.group(2), float(m2.group(1)))

    # If not enough from key_points, mine the description
    if len(results) < 2 and description:
        for m in pct_inline.finditer(description):
            _add(m.group(2), float(m.group(1)))
        if len(results) < 2:
            for m in pct_labeled.finditer(description):
                _add(m.group(1), float(m.group(2)))

    return results if len(results) >= 2 else []

def extract_bar_data(key_points: List[str], description: str = ""):
    """
    Extract (label, value) pairs from key_points AND description that contain
    comparable numeric values (not percentages).
    e.g. "Revenue: 4.5 billion" -> ("Revenue", 4.5)
    """
    results = []
    seen_labels = set()

    # Strict pattern: requires a clean label (1-4 words) before the number
    num_pattern = _re.compile(
        r'\b([A-Z][a-zA-Z\s&/\(\)]{1,30}?)[\s:\-–]+\$?(\d+(?:\.\d+)?)\s*(million|billion|trillion|thousand|k|m|b)?\b(?!\s*%)',
        _re.IGNORECASE
    )

    def _add(label, value, multiplier):
        label = label.strip().strip('•:–-').strip()
        # Reject labels that are too short, too long, or look like sentence fragments
        words = label.split()
        if not label or len(words) > 4 or len(words) < 1:
            return
        if label.lower() in seen_labels:
            return
        mult = (multiplier or '').lower()
        if mult in ('trillion',):          value *= 1_000_000_000_000
        elif mult in ('billion', 'b'):     value *= 1_000_000_000
        elif mult in ('million', 'm'):     value *= 1_000_000
        elif mult in ('thousand', 'k'):    value *= 1_000
        seen_labels.add(label.lower())
        results.append((label[:40], value))

    # Key points first (most reliable)
    for pt in key_points:
        m = num_pattern.match(pt.strip().strip('•').strip())
        if m:
            _add(m.group(1), float(m.group(2)), m.group(3))

    # Description fallback — only scan bullet-like sentences (short, colon-separated)
    if len(results) < 2 and description:
        # Only match "Label: value unit" style in description, not mid-sentence
        strict = _re.compile(
            r'(?:^|[,\.])\s*([A-Z][a-zA-Z\s]{1,25}?):\s*\$?(\d+(?:\.\d+)?)\s*(million|billion|trillion|thousand)?\b(?!\s*%)',
            _re.IGNORECASE | _re.MULTILINE
        )
        for m in strict.finditer(description):
            _add(m.group(1), float(m.group(2)), m.group(3))

    return results if len(results) >= 2 else []

def detect_chart_type(slide: Dict) -> str:
    """
    Returns 'pie', 'bar', 'line', 'area', 'scatter', or '' based on content.
    Pie: 2+ percentage values present.
    Bar: 2+ comparable numeric values (non-percentage).
    Line/Area: Time series data or trend keywords present.
    Scatter: Multiple categories with coordinate pairs.
    """
    key_points  = slide.get('key_points', [])
    description = slide.get('description', '')
    title = slide.get('title', '').lower()
    all_text = (title + " " + description).lower()
    
    if len(extract_pie_data(key_points, description)) >= 2:
        return 'pie'
    
    # Check for time series/trend keywords
    trend_keywords = ['trend', 'growth', 'over time', 'year', 'quarter', 'month', 'forecast', 'projection', 'historical', 'timeline']
    if any(kw in all_text for kw in trend_keywords):
        time_labels, time_values = extract_time_series_data(key_points, description)
        if len(time_labels) >= 3:
            if any(kw in all_text for kw in ['forecast', 'projection', 'future']):
                return 'area'  # Area for composition/accumulation
            return 'line'  # Line for trends over time
    
    # Check for multi-series comparison
    series_data = extract_comparison_series(key_points, description)
    if len(series_data) >= 2:
        return 'bar'  # Multi-series bar chart
    
    if len(extract_bar_data(key_points, description)) >= 2:
        return 'bar'
    
    return ''


def detect_slide_layout(slide: Dict, is_last_slide: bool = False, has_image: bool = False) -> str:
    """
    Intelligently detect the best layout for a slide based on its content.
    Returns one of: standard, full_image, large_key_points, stats, timeline, pie_chart, bar_chart, line_chart, area_chart
    """
    title      = slide.get('title', '').lower()
    description = slide.get('description', '').lower()
    key_points  = slide.get('key_points', [])

    # Priority 1: Last slide gets large_key_points layout
    if is_last_slide and key_points and len(key_points) <= 5:
        if any(kw in title for kw in ['conclusion', 'summary', 'key takeaways', 'final', 'wrap', 'future']):
            return "large_key_points"

    # Priority 1a: Chart layouts — only when no image and not last slide
    if not has_image and not is_last_slide:
        chart = detect_chart_type(slide)
        if chart == 'pie':
            return 'pie_chart'
        if chart == 'bar':
            return 'bar_chart'
        if chart == 'line':
            return 'line_chart'
        if chart == 'area':
            return 'area_chart'

    # Priority 1b: Stats layout for numeric content (no image needed)
    if not has_image and has_numeric_content(slide):
        if not is_last_slide:
            return "stats"

    # Priority 1c: Timeline layout for process/steps slides
    TIMELINE_KEYWORDS = ['step', 'phase', 'stage', 'process', 'how it works', 'workflow',
                         'pipeline', 'roadmap', 'journey', 'evolution', 'history', 'timeline']
    if not has_image and not is_last_slide:
        if any(kw in title for kw in TIMELINE_KEYWORDS):
            return "timeline"

    # Priority 2: Full image layout if image available and short description
    if has_image and len(description) < 150:
        return "full_image"

    # Priority 3: Standard layout with image if available
    if has_image:
        return "standard"

    return "standard"

# THEME COLOR SCHEMES 
THEMES = {
    "modern": {
        "primary": RGBColor(30, 64, 175),      # #1e40af (blue)
        "primary_light": RGBColor(59, 130, 246),  # #3b82f6
        "accent": RGBColor(245, 158, 11),     # #f59e0b (amber)
        "white": RGBColor(255, 255, 255),
        "bg_light": RGBColor(248, 250, 252),  # #f8fafc
        "text_primary": RGBColor(15, 23, 42),  # #0f172a
        "text_secondary": RGBColor(56, 65, 83),  # #384153 (stronger contrast)
        "border": RGBColor(226, 232, 240),    # #e2e8f0
    },
    "dark": {
        "primary": RGBColor(17, 24, 39),      # Dark gray
        "primary_light": RGBColor(31, 41, 55),
        "accent": RGBColor(34, 197, 94),      # Bright green
        "white": RGBColor(255, 255, 255),
        "bg_light": RGBColor(20, 28, 40),     # Very dark
        "text_primary": RGBColor(243, 244, 246),  # Light gray
        "text_secondary": RGBColor(209, 213, 219),
        "border": RGBColor(75, 85, 99),
    },
    "minimal": {
        "primary": RGBColor(55, 65, 81),      # Gray
        "primary_light": RGBColor(107, 114, 128),
        "accent": RGBColor(59, 130, 246),     # Light blue
        "white": RGBColor(255, 255, 255),
        "bg_light": RGBColor(249, 250, 251),  # Very light gray
        "text_primary": RGBColor(17, 24, 39),  # Dark gray
        "text_secondary": RGBColor(107, 114, 128),
        "border": RGBColor(229, 231, 235),
    }
}

class PresentationGenerator:
    def __init__(self, title: str, subtitle: str = "", theme: str = "modern"):
        self.prs = PPTXPresentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle
        self.theme = theme if theme in THEMES else "modern"
        self.colors = THEMES[self.theme]
        
    def add_title_slide(self, logo_path: str = None):
        """Add title slide matching the indigo design with decorative circles."""
        from pptx.oxml.ns import qn
        from lxml import etree

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # ── Background: theme-aware color ──────────────────────────────────────
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["primary"]

        # ── Large filled circle – top-right (partially off-slide) ────────────
        # Diameter ~4.5 in, centre near (9.2, 1.0)
        big_r = 2.25
        big_circle = slide.shapes.add_shape(
            9,  # oval
            Inches(9.2 - big_r), Inches(1.0 - big_r),
            Inches(big_r * 2),   Inches(big_r * 2)
        )
        big_circle.fill.solid()
        big_circle.fill.fore_color.rgb = self.colors["primary_light"]
        big_circle.line.fill.background()

        # ── Small outline circle – bottom-left ───────────────────────────────
        sm_r = 1.10
        sm_circle = slide.shapes.add_shape(
            9,
            Inches(0.10), Inches(7.5 - sm_r * 2 - 0.10),
            Inches(sm_r * 2), Inches(sm_r * 2)
        )
        sm_circle.fill.background()                               # transparent fill
        sm_circle.line.color.rgb = self.colors["primary_light"]
        sm_circle.line.width = Pt(2.5)

        # ── Title text – centred vertically and horizontally ─────────────────
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(2.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.15

        # ── Optional subtitle ─────────────────────────────────────────────────
        if self.subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.8))
            sf = sub_box.text_frame
            sp = sf.paragraphs[0]
            sp.text = self.subtitle
            sp.font.size = Pt(20)
            sp.font.color.rgb = RGBColor(197, 202, 233)   # light indigo tint
            sp.alignment = PP_ALIGN.CENTER

        # ── Thin accent bar at the very bottom ───────────────────────────────
        bar = slide.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(10), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(92, 107, 192)
        bar.line.fill.background()
    
    def _add_footer(self, slide, page_num: int):
        """Add a consistent footer with topic name + page number to any content slide."""
        # Thin separator line
        sep = slide.shapes.add_shape(1, Inches(0.5), Inches(7.15), Inches(9), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self.colors["border"]
        sep.line.fill.background()

        # Topic name (left)
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(7), Inches(0.28))
        lf = left_box.text_frame
        lp = lf.paragraphs[0]
        lp.text = self.title
        lp.font.size = Pt(9)
        lp.font.color.rgb = self.colors["text_secondary"]

        # Page number (right)
        right_box = slide.shapes.add_textbox(Inches(8.5), Inches(7.18), Inches(1), Inches(0.28))
        rf = right_box.text_frame
        rp = rf.paragraphs[0]
        rp.text = str(page_num)
        rp.font.size = Pt(9)
        rp.font.color.rgb = self.colors["text_secondary"]
        rp.alignment = PP_ALIGN.RIGHT

    def _add_gradient_bg(self, slide):
        """
        Solid white background for modern theme.
        Subtle tint only for dark/minimal themes.
        """
        if self.theme == "dark":
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = self.colors["bg_light"]
        else:
            # modern + minimal: pure white, no tint
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = self.colors["white"]

    def _heading_color(self):
        """Return a readable heading color for the active theme."""
        return self.colors["white"] if self.theme == "dark" else self.colors["primary"]

    def _body_color(self):
        """Return the main body text color for the active theme."""
        return self.colors["text_primary"]

    def add_section_divider(self, section_title: str, section_subtitle: str = "", page_num: int = 0):
        """
        Full-bleed section divider slide — primary color background,
        large section number on the right, title centered left.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["primary"]

        # Large decorative circle top-right
        circle = slide.shapes.add_shape(9, Inches(7.5), Inches(-1.0), Inches(4.5), Inches(4.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.colors["primary_light"]
        circle.line.fill.background()

        # Accent bar left edge
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(8.5), Inches(1.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = self.colors["white"]
        p.line_spacing = 1.1

        # Optional subtitle
        if section_subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(8), Inches(0.6))
            sf = sub_box.text_frame
            sp = sf.paragraphs[0]
            sp.text = section_subtitle
            sp.font.size = Pt(16)
            sp.font.color.rgb = RGBColor(197, 202, 233)

        # Bottom accent line
        bottom = slide.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(10), Inches(0.12))
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = self.colors["accent"]
        bottom.line.fill.background()

    def add_content_slide(self, title: str, content: List[str], slide_type: str = "bullet", page_num: int = 0):
        """Add content slide — gradient bg, icon circles for bullets, footer."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # ── Gradient background (two-layer simulation) ────────────────────────
        self._add_gradient_bg(slide)

        # ── Header accent line ────────────────────────────────────────────────
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors["accent"]
        header_shape.line.color.rgb = self.colors["accent"]

        # ── Title ─────────────────────────────────────────────────────────────
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["white"] if self.theme == "dark" else self.colors["primary"]

        # ── Content ───────────────────────────────────────────────────────────
        if slide_type == "bullet":
            # Split into para lines and bullet lines
            para_lines  = [c for c in content if c and not c.startswith("•")]
            bullet_lines = [c for c in content if c.startswith("•")]

            # Two-column layout when 3-4 bullets (no para) — otherwise single column
            use_two_col = len(bullet_lines) >= 3 and not para_lines

            if use_two_col:
                # Split bullets into two columns
                mid = (len(bullet_lines) + 1) // 2
                col_bullets = [bullet_lines[:mid], bullet_lines[mid:]]
                col_x = [Inches(0.6), Inches(5.2)]
                col_w = Inches(4.4)

                for col_idx, (cx, col_items) in enumerate(zip(col_x, col_bullets)):
                    cb = slide.shapes.add_textbox(cx, Inches(1.55), col_w, Inches(5.3))
                    ctf = cb.text_frame
                    ctf.word_wrap = True
                    for bi, item in enumerate(col_items):
                        bp = ctf.paragraphs[0] if bi == 0 else ctf.add_paragraph()
                        self._render_bullet(slide, bp, item, cx, Inches(1.55 + bi * 0.85))
            else:
                # Single column — para first, then bullets with icon circles
                y_cursor = 1.55  # inches from top

                # Render para lines as plain text
                if para_lines:
                    para_text = " ".join(para_lines)
                    # More conservative estimate: ~70 chars per line at Pt(16) / 8.4 inch width
                    chars_per_line = 70
                    estimated_lines = max(1, -(-len(para_text) // chars_per_line))  # ceiling div
                    para_height = max(0.50, estimated_lines * 0.28 + 0.10)

                    pb = slide.shapes.add_textbox(Inches(0.8), Inches(y_cursor), Inches(8.4), Inches(para_height))
                    ptf = pb.text_frame
                    ptf.word_wrap = True
                    pp = ptf.paragraphs[0]
                    pp.text = para_text
                    pp.font.size = Pt(16)
                    pp.font.color.rgb = self.colors["text_primary"]
                    pp.line_spacing = 1.3
                    y_cursor += para_height + 0.22  # tighter gap so bullets have more room

                # Render each bullet with an icon circle
                # Footer is at y=7.15, leave 0.5 inch margin → max bullet bottom = 6.65
                MAX_Y = 6.60
                for item in bullet_lines:
                    if y_cursor >= MAX_Y:
                        break  # don't draw bullets that would overlap the footer
                    self._add_bullet_row(slide, item, y_cursor)
                    y_cursor += 0.76
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(5.0))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = "\n".join(content)
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors["text_primary"]
            p.line_spacing = 1.5

        # ── Footer ────────────────────────────────────────────────────────────
        if page_num:
            self._add_footer(slide, page_num)

    def _render_bullet(self, slide, paragraph, text: str, x_inches, y_inches):
        """Render a bullet paragraph (used in two-column mode — no icon circle)."""
        paragraph.text = text  # keeps the • character
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = self.colors["text_primary"]
        paragraph.space_before = Pt(8)
        paragraph.space_after  = Pt(8)
        paragraph.line_spacing = 1.3

    def _add_bullet_row(self, slide, text: str, y_top: float):
        """
        Render one bullet row: small filled circle on the left + text to the right.
        y_top is in inches from the top of the slide.
        """
        CIRCLE_D = 0.18   # diameter of icon circle in inches
        TEXT_X   = 0.95   # where text starts
        TEXT_W   = 8.3

        # Icon circle
        circle = slide.shapes.add_shape(
            9,  # oval
            Inches(0.55), Inches(y_top + 0.05),
            Inches(CIRCLE_D), Inches(CIRCLE_D)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.colors["accent"]
        circle.line.fill.background()

        # Text
        tb = slide.shapes.add_textbox(Inches(TEXT_X), Inches(y_top), Inches(TEXT_W), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Strip the leading "• " since we have the circle
        p.text = text.lstrip("• ").strip()
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors["text_primary"]
        p.line_spacing = 1.2
    
    def add_full_image_slide(self, title: str, image_path: str, description: str = ""):
        """Add full-width image slide with optional text overlay"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["white"]
        
        # Add image to fill slide
        try:
            slide.shapes.add_picture(image_path, Inches(0), Inches(1.2), width=Inches(10), height=Inches(5.3))
        except Exception as e:
            print(f"WARNING: Could not add image {image_path}: {e}")
        
        # Title bar (top)
        title_bg = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.9))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors["primary"]
        title_bg.line.color.rgb = self.colors["primary"]
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER
        
        if description:
            desc_bg = slide.shapes.add_shape(1, Inches(0), Inches(6), Inches(10), Inches(1.5))
            desc_bg.fill.solid()
            desc_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
            desc_bg.fill.transparency = 0.3
            desc_bg.line.color.rgb = RGBColor(0, 0, 0)
            
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9), Inches(1.3))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            p = desc_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors["white"]
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.3
    
    def add_quote_slide(self, quote: str, author: str = ""):
        """Add centered quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]
        
        # Left accent bar
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors["accent"]
        accent_bar.line.color.rgb = self.colors["accent"]
        
        # Quote text (centered)
        quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(2.5))
        quote_frame = quote_box.text_frame
        quote_frame.word_wrap = True
        quote_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = quote_frame.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.4
        
        # Author (if provided)
        if author:
            author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(0.8))
            author_frame = author_box.text_frame
            ap = author_frame.paragraphs[0]
            ap.text = f"— {author}"
            ap.font.size = Pt(18)
            ap.font.color.rgb = self.colors["accent"]
            ap.alignment = PP_ALIGN.CENTER
    
    def add_large_key_points_slide(self, title: str, key_points: List[str]):
        """Add slide with 3-5 large key points (1 per line)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["white"] if self.theme != "dark" else self.colors["bg_light"]
        
        # Header with accent line
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors["accent"]
        header_shape.line.color.rgb = self.colors["accent"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self._heading_color()
        
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(7), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(key_points[:5]): 
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{i + 1}. {point}"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.colors["text_primary"] if self.theme == "dark" else self.colors["primary"]
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.line_spacing = 1.3
    
    def add_comparison_slide(self, title: str, items: List[Dict]):
        """
        Staggered diagonal flow slide.

        Three cards are placed at descending vertical positions (top → middle → bottom)
        with a diagonal connector line running through their centre dots, giving a clear
        sense of progression.  Each card has:
          • A filled circle with the step number on the left edge
          • A bold title
          • Supporting bullet points below the title

        Layout (slide is 10 × 7.5 inches):
          Card 1  x=0.55  y=1.20  (top)
          Card 2  x=3.30  y=2.70  (middle)
          Card 3  x=6.05  y=4.20  (bottom)
          Card width=3.40  height=2.10
        """
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        from lxml import etree

        CARD_W   = 3.40   # inches
        CARD_H   = 2.10
        CIRCLE_R = 0.28   # radius of the step-number circle
        MARGIN   = 0.18   # inner padding inside card

        # Horizontal anchors for each card (left edge)
        card_x = [0.55, 3.30, 6.05]
        # Vertical anchors – staggered top → middle → bottom
        card_y = [1.20, 2.70, 4.20]

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = (
            self.colors["bg_light"] if self.theme != "dark" else self.colors["primary"]
        )

        # ── Thin accent bar at very top ──────────────────────────────────────
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.28), Inches(10), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.color.rgb     = self.colors["accent"]

        # ── Slide title ──────────────────────────────────────────────────────
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.36), Inches(9), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size  = Pt(26)
        p.font.bold  = True
        p.font.color.rgb = self._heading_color()

        # ── Helper: centre-point of each card's circle ───────────────────────
        # Circle sits on the left edge of the card, vertically centred at top 1/3
        def circle_centre(cx, cy):
            """Return (x, y) inches for the centre of the step circle."""
            return (cx + CIRCLE_R, cy + CIRCLE_R + 0.10)

        # ── Diagonal connector line (drawn first so cards sit on top) ────────
        # We draw one straight line from circle-1 centre to circle-3 centre.
        centres = [circle_centre(card_x[i], card_y[i]) for i in range(len(items[:3]))]

        if len(centres) >= 2:
            x1, y1 = centres[0]
            x2, y2 = centres[-1]

            # python-pptx doesn't expose connectors directly, so we use a
            # thin rectangle rotated to act as a line.  A simpler approach:
            # add a connector shape via the XML directly.
            sp_tree = slide.shapes._spTree

            def emu(inches):
                return int(inches * 914400)

            cxn_xml = (
                '<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                '<p:nvCxnSpPr>'
                '<p:cNvPr id="900" name="DiagConnector"/>'
                '<p:cNvCxnSpPr/>'
                '<p:nvPr/>'
                '</p:nvCxnSpPr>'
                '<p:spPr>'
                f'<a:xfrm><a:off x="{emu(x1)}" y="{emu(y1)}"/>'
                f'<a:ext cx="{emu(x2 - x1)}" cy="{emu(y2 - y1)}"/></a:xfrm>'
                '<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
                '<a:ln w="18000">'          # 2 pt line
                '<a:solidFill><a:srgbClr val="CBD5E1"/></a:solidFill>'
                '<a:prstDash val="sysDash"/>'
                '</a:ln>'
                '</p:spPr>'
                '</p:cxnSp>'
            )
            sp_tree.append(etree.fromstring(cxn_xml))

        # ── Draw each card ───────────────────────────────────────────────────
        for idx, item in enumerate(items[:3]):
            cx = card_x[idx]
            cy = card_y[idx]
            cc_x, cc_y = circle_centre(cx, cy)

            # Card shadow strip (offset 0.06 right + down, slightly darker bg)
            shadow = slide.shapes.add_shape(
                1,
                Inches(cx + 0.06), Inches(cy + 0.06),
                Inches(CARD_W),    Inches(CARD_H)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = self.colors["border"]
            shadow.line.fill.background()   # no border on shadow

            # Card body
            card = slide.shapes.add_shape(
                1,
                Inches(cx), Inches(cy),
                Inches(CARD_W), Inches(CARD_H)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = (
                self.colors["white"] if self.theme != "dark" else self.colors["primary_light"]
            )
            card.line.color.rgb = self.colors["border"]
            card.line.width     = Pt(0.75)

            # Accent left border strip
            accent_strip = slide.shapes.add_shape(
                1,
                Inches(cx), Inches(cy),
                Inches(0.06), Inches(CARD_H)
            )
            accent_strip.fill.solid()
            accent_strip.fill.fore_color.rgb = self.colors["accent"]
            accent_strip.line.fill.background()

            # Step-number circle
            circle = slide.shapes.add_shape(
                9,   # MSO_SHAPE_TYPE.OVAL
                Inches(cc_x - CIRCLE_R), Inches(cc_y - CIRCLE_R),
                Inches(CIRCLE_R * 2),    Inches(CIRCLE_R * 2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors["primary"]
            circle.line.color.rgb      = self.colors["primary"]

            num_tf = circle.text_frame
            num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np_ = num_tf.paragraphs[0]
            np_.text = str(idx + 1)
            np_.font.size  = Pt(13)
            np_.font.bold  = True
            np_.font.color.rgb = self.colors["white"]
            np_.alignment  = PP_ALIGN.CENTER

            # Card title  (starts after the circle, with left padding)
            title_left  = cx + CIRCLE_R * 2 + MARGIN + 0.10
            title_width = CARD_W - (CIRCLE_R * 2 + MARGIN + 0.10) - MARGIN
            title_top   = cy + MARGIN

            t_box = slide.shapes.add_textbox(
                Inches(title_left), Inches(title_top),
                Inches(title_width), Inches(0.42)
            )
            t_tf = t_box.text_frame
            t_tf.word_wrap = True
            tp = t_tf.paragraphs[0]
            tp.text = item.get("title", f"Step {idx + 1}")
            tp.font.size  = Pt(14)
            tp.font.bold  = True
            tp.font.color.rgb = self.colors["text_primary"] if self.theme == "dark" else self.colors["primary"]

            # Thin separator line under title
            sep = slide.shapes.add_shape(
                1,
                Inches(cx + 0.18), Inches(cy + 0.66),
                Inches(CARD_W - 0.36), Inches(0.02)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = self.colors["border"]
            sep.line.fill.background()

            # Bullet points
            raw_content = item.get("content", "")
            bullets = [b.strip() for b in raw_content.split("\n") if b.strip()][:3]

            c_box = slide.shapes.add_textbox(
                Inches(cx + 0.22), Inches(cy + 0.74),
                Inches(CARD_W - 0.40), Inches(CARD_H - 0.80)
            )
            c_tf = c_box.text_frame
            c_tf.word_wrap = True

            for bi, bullet in enumerate(bullets):
                bp = c_tf.paragraphs[0] if bi == 0 else c_tf.add_paragraph()
                bp.text = f"• {bullet}"
                bp.font.size  = Pt(11)
                bp.font.color.rgb = self._body_color()
                bp.space_before   = Pt(3)
                bp.space_after    = Pt(3)
                bp.line_spacing   = 1.2
    
    def add_summary_slide(self, key_points: List[str]):
        """Key Takeaways slide — numbered white cards on primary background."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["primary"]

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.30), Inches(9), Inches(0.65))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Takeaways"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER

        # Accent underline bar
        ubar = slide.shapes.add_shape(1, Inches(3.5), Inches(0.98), Inches(3.0), Inches(0.05))
        ubar.fill.solid()
        ubar.fill.fore_color.rgb = self.colors["accent"]
        ubar.line.fill.background()

        CARD_H   = 0.85
        CARD_GAP = 0.12
        CARD_X   = 1.0
        CARD_W   = 8.0
        CIRCLE_D = 0.32
        START_Y  = 1.15

        for i, point in enumerate(key_points[:5]):
            cy = START_Y + i * (CARD_H + CARD_GAP)

            # Card shadow
            shadow = slide.shapes.add_shape(1, Inches(CARD_X + 0.05), Inches(cy + 0.05), Inches(CARD_W), Inches(CARD_H))
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = self.colors["primary_light"]
            shadow.line.fill.background()

            # Card body
            card = slide.shapes.add_shape(1, Inches(CARD_X), Inches(cy), Inches(CARD_W), Inches(CARD_H))
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors["white"]
            card.line.color.rgb = self.colors["border"]
            card.line.width = Pt(0.5)

            # Number circle
            circle_x = CARD_X + 0.22
            circle_y = cy + (CARD_H - CIRCLE_D) / 2
            circle = slide.shapes.add_shape(9, Inches(circle_x), Inches(circle_y), Inches(CIRCLE_D), Inches(CIRCLE_D))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors["accent"]
            circle.line.fill.background()
            ntf = circle.text_frame
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np_ = ntf.paragraphs[0]
            np_.text = str(i + 1)
            np_.font.size = Pt(12)
            np_.font.bold = True
            np_.font.color.rgb = self.colors["white"]
            np_.alignment = PP_ALIGN.CENTER

            # Takeaway text
            text_x = CARD_X + CIRCLE_D + 0.38
            text_w = CARD_W - CIRCLE_D - 0.50
            ttb = slide.shapes.add_textbox(Inches(text_x), Inches(cy + 0.18), Inches(text_w), Inches(0.50))
            ttf = ttb.text_frame
            ttf.word_wrap = True
            tp = ttf.paragraphs[0]
            tp.text = point
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = self.colors["primary"]
            tp.line_spacing = 1.2
    
    def add_enhanced_conclusion_slide(self, title: str, key_points: List[str]):
        """Add enhanced 3-column conclusion slide: What We Learned | Key Actions | What's Next"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["white"] if self.theme != "dark" else self.colors["bg_light"]
        
        # Accent bar at top
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors["accent"]
        accent_bar.line.color.rgb = self.colors["accent"]
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._heading_color()
        p.alignment = PP_ALIGN.CENTER
        
        # Column widths and positions
        col_width = 2.8
        col_height = 3.8
        col_positions = [0.8, 3.8, 6.8]
        column_headers = ["What We\nLearned", "Key\nActions", "What's\nNext"]
        
        # Divider lines between columns
        for x in [3.5, 6.5]:
            divider = slide.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(0.02), Inches(4.2))
            divider.fill.solid()
            divider.fill.fore_color.rgb = self.colors["border"]
            divider.line.color.rgb = self.colors["border"]
        
        # Create 3 columns
        for col_idx, (x_pos, header) in enumerate(zip(col_positions, column_headers)):
            # Column header
            header_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(col_width), Inches(0.6))
            header_frame = header_box.text_frame
            header_frame.word_wrap = True
            p = header_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors["accent"]
            p.alignment = PP_ALIGN.CENTER
            
            # Column content
            content_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(2.1), Inches(col_width - 0.2), Inches(col_height))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            # Distribute key_points across 3 columns
            start_idx = col_idx * (len(key_points) // 3)
            end_idx = start_idx + (len(key_points) // 3)
            if col_idx == 2:  # Last column gets remainder
                end_idx = len(key_points)
            
            column_points = key_points[start_idx:end_idx]
            
            for i, point in enumerate(column_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = "• " + point
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors["text_primary"]
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                p.line_spacing = 1.2
        
        # Footer section with branding
        footer_line = slide.shapes.add_shape(1, Inches(0.5), Inches(5.6), Inches(9), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.colors["border"]
        footer_line.line.color.rgb = self.colors["border"]
        
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        p = footer_frame.paragraphs[0]
        p.text = "Thank you for your attention | Questions?"
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors["text_secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide_with_image(self, title: str, content: List[str], image_path: str, image_position: str = "right", page_num: int = 0):
        """Add content slide with image in 2-column layout — gradient bg + footer."""
        import os

        if not os.path.exists(image_path):
            print(f"DEBUG: Image not found at {image_path}, adding regular slide instead")
            self.add_content_slide(title, content, slide_type='bullet', page_num=page_num)
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # ── Gradient background ───────────────────────────────────────────────
        self._add_gradient_bg(slide)

        # ── Header accent line ────────────────────────────────────────────────
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors["accent"]
        header_shape.line.color.rgb = self.colors["accent"]

        # ── Title ─────────────────────────────────────────────────────────────
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self._heading_color()

        if image_position == "right":
            text_x    = Inches(0.5)
            text_width = Inches(4.5)
            image_x   = Inches(5.2)
        else:
            image_x   = Inches(0.5)
            text_x    = Inches(5.2)
            text_width = Inches(4.3)

        # ── Text content (para + icon-circle bullets) ─────────────────────────
        para_lines   = [c for c in content if c and not c.startswith("•")]
        bullet_lines = [c for c in content if c.startswith("•")]

        # text_x is an Inches object — extract float value for arithmetic
        tx = text_x.inches if hasattr(text_x, 'inches') else 0.5
        tw = text_width.inches if hasattr(text_width, 'inches') else 4.5

        y_cursor = 1.55
        if para_lines:
            para_text = " ".join(para_lines)
            # 4.5 inch wide column at Pt(14) ≈ 55 chars per line
            chars_per_line = 55
            estimated_lines = max(1, -(-len(para_text) // chars_per_line))  # ceiling div
            para_height = max(0.50, estimated_lines * 0.32 + 0.15)

            pb = slide.shapes.add_textbox(Inches(tx), Inches(y_cursor), Inches(tw), Inches(para_height))
            ptf = pb.text_frame
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = para_text
            pp.font.size = Pt(14)
            pp.font.color.rgb = self.colors["text_primary"]
            pp.line_spacing = 1.3
            y_cursor += para_height + 0.35  # clear gap before first bullet

        for item in bullet_lines:
            # Icon circle
            circle = slide.shapes.add_shape(
                9, Inches(tx), Inches(y_cursor + 0.06), Inches(0.16), Inches(0.16)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors["accent"]
            circle.line.fill.background()

            tb = slide.shapes.add_textbox(
                Inches(tx + 0.26), Inches(y_cursor),
                Inches(tw - 0.28), Inches(0.72)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            bp = tf.paragraphs[0]
            bp.text = item.lstrip("• ").strip()
            bp.font.size = Pt(15)
            bp.font.color.rgb = self.colors["text_primary"]
            bp.line_spacing = 1.3
            y_cursor += 0.80

        # ── Image ─────────────────────────────────────────────────────────────
        try:
            slide.shapes.add_picture(image_path, image_x, Inches(1.4), width=Inches(4.2), height=Inches(5.4))
            print(f"DEBUG: Image added to slide with title '{title}'")
        except Exception as e:
            print(f"WARNING: Could not add image to slide: {str(e)}")

        # ── Footer ────────────────────────────────────────────────────────────
        if page_num:
            self._add_footer(slide, page_num)
    
    
    def add_stats_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Stats highlight slide — extract numbers and render as large callout cards.
        Non-numeric points fall back to regular bullet rows below the cards.
        """
        import re

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Split points into numeric (for cards) and text (for bullets)
        numeric_points = []
        text_points = []
        num_pattern = re.compile(r'(\d+%|\$[\d,]+|\d{4}|\d+\s*(?:million|billion|thousand|x|times)|\b\d{2,}\b)', re.IGNORECASE)

        for pt in key_points[:6]:
            m = num_pattern.search(pt)
            if m and len(numeric_points) < 3:
                number = m.group(0)
                label = pt.replace(number, '').strip().strip(':-–').strip()
                if not label:
                    label = pt
                numeric_points.append((number, label))
            else:
                text_points.append(pt)

        # Stat cards row
        if numeric_points:
            n = len(numeric_points)
            card_w = min(2.8, 8.6 / n)
            total_w = n * card_w + (n - 1) * 0.2
            start_x = (10 - total_w) / 2

            for ci, (number, label) in enumerate(numeric_points):
                cx = start_x + ci * (card_w + 0.2)
                cy = 1.35

                # Card
                card = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(card_w), Inches(1.9))
                card.fill.solid()
                card.fill.fore_color.rgb = self.colors["bg_light"]
                card.line.color.rgb = self.colors["border"]
                card.line.width = Pt(0.75)

                # Top accent strip
                strip = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(card_w), Inches(0.07))
                strip.fill.solid()
                strip.fill.fore_color.rgb = self.colors["accent"]
                strip.line.fill.background()

                # Big number
                nb = slide.shapes.add_textbox(Inches(cx + 0.1), Inches(cy + 0.12), Inches(card_w - 0.2), Inches(0.9))
                ntf = nb.text_frame
                np_ = ntf.paragraphs[0]
                np_.text = number
                np_.font.size = Pt(36)
                np_.font.bold = True
                np_.font.color.rgb = self.colors["primary"]
                np_.alignment = PP_ALIGN.CENTER

                # Label
                lb = slide.shapes.add_textbox(Inches(cx + 0.1), Inches(cy + 1.05), Inches(card_w - 0.2), Inches(0.75))
                ltf = lb.text_frame
                ltf.word_wrap = True
                lp = ltf.paragraphs[0]
                lp.text = label
                lp.font.size = Pt(12)
                lp.font.color.rgb = self._body_color()
                lp.alignment = PP_ALIGN.CENTER
                lp.line_spacing = 1.2

        # Description para
        y_cursor = 3.45 if numeric_points else 1.55
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            chars_per_line = 95
            estimated_lines = max(1, -(-len(short_desc) // chars_per_line))
            para_height = max(0.45, estimated_lines * 0.30 + 0.10)
            pb = slide.shapes.add_textbox(Inches(0.8), Inches(y_cursor), Inches(8.4), Inches(para_height))
            ptf = pb.text_frame
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(14)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_cursor += para_height + 0.25

        # Remaining text bullets
        for pt in text_points[:3]:
            self._add_bullet_row(slide, f"• {pt}", y_cursor)
            y_cursor += 0.80

        if page_num:
            self._add_footer(slide, page_num)

    def add_timeline_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Horizontal timeline slide — each key point is a milestone on a line.
        Works best with 3-5 points.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Description
        y_after_desc = 1.45
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            pb = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(8.4), Inches(0.55))
            ptf = pb.text_frame
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(13)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_after_desc = 2.10

        points = key_points[:5]
        n = len(points)
        if n == 0:
            return

        # Timeline horizontal line
        LINE_Y  = y_after_desc + 1.10   # y of the centre line
        LINE_X1 = 0.8
        LINE_X2 = 9.2
        line_bar = slide.shapes.add_shape(1, Inches(LINE_X1), Inches(LINE_Y - 0.03), Inches(LINE_X2 - LINE_X1), Inches(0.06))
        line_bar.fill.solid()
        line_bar.fill.fore_color.rgb = self.colors["border"]
        line_bar.line.fill.background()

        # Evenly space milestone dots
        spacing = (LINE_X2 - LINE_X1) / (n - 1) if n > 1 else 0
        DOT_R = 0.18

        for idx, point in enumerate(points):
            dot_x = LINE_X1 + idx * spacing
            dot_y = LINE_Y

            # Dot
            dot = slide.shapes.add_shape(9, Inches(dot_x - DOT_R), Inches(dot_y - DOT_R), Inches(DOT_R * 2), Inches(DOT_R * 2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.colors["accent"] if idx % 2 == 0 else self.colors["primary"]
            dot.line.color.rgb = self.colors["white"]
            dot.line.width = Pt(1.5)

            # Step number inside dot
            dtf = dot.text_frame
            dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
            dp = dtf.paragraphs[0]
            dp.text = str(idx + 1)
            dp.font.size = Pt(10)
            dp.font.bold = True
            dp.font.color.rgb = self.colors["white"]
            dp.alignment = PP_ALIGN.CENTER

            # Label box — alternate above/below the line for readability
            label_w = min(2.2, spacing + 0.2) if n > 1 else 2.2
            label_x = dot_x - label_w / 2

            if idx % 2 == 0:
                # Above the line
                label_y = dot_y - DOT_R - 1.10
                label_h = 0.90
            else:
                # Below the line
                label_y = dot_y + DOT_R + 0.15
                label_h = 0.90

            ltb = slide.shapes.add_textbox(Inches(label_x), Inches(label_y), Inches(label_w), Inches(label_h))
            ltf = ltb.text_frame
            ltf.word_wrap = True
            lp = ltf.paragraphs[0]
            lp.text = point
            lp.font.size = Pt(12)
            lp.font.bold = (idx % 2 == 0)
            lp.font.color.rgb = self.colors["primary"] if idx % 2 == 0 else self._body_color()
            lp.alignment = PP_ALIGN.CENTER
            lp.line_spacing = 1.2

            # Connector tick line from dot to label
            tick_x = dot_x - 0.01
            if idx % 2 == 0:
                tick_y = label_y + label_h
                tick_h = dot_y - DOT_R - (label_y + label_h)
            else:
                tick_y = dot_y + DOT_R
                tick_h = label_y - (dot_y + DOT_R)

            if tick_h > 0.05:
                tick = slide.shapes.add_shape(1, Inches(tick_x), Inches(tick_y), Inches(0.02), Inches(tick_h))
                tick.fill.solid()
                tick.fill.fore_color.rgb = self.colors["border"]
                tick.line.fill.background()

        if page_num:
            self._add_footer(slide, page_num)

    def add_pie_chart_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Split layout: description + non-chart bullets on the left, pie chart on the right.
        Chart data is extracted from percentage key_points.
        """
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        import re

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Extract chart data
        pie_data = extract_pie_data(key_points, description)
        text_points = [pt for pt in key_points
                       if not re.search(r'\d+(?:\.\d+)?\s*%', pt)]

        # Left column — description + non-chart bullets
        y_cursor = 1.45
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            # Conservative estimate for lines to avoid clipping on narrow textboxes
            chars_per_line = 40
            est_lines = max(1, -(-len(short_desc) // chars_per_line))
            ph = max(0.55, est_lines * 0.34 + 0.20)
            pb = slide.shapes.add_textbox(Inches(0.5), Inches(y_cursor), Inches(4.3), Inches(ph))
            ptf = pb.text_frame
            try:
                ptf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(13)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_cursor += ph + 0.30

        for pt in text_points[:4]:
            self._add_bullet_row(slide, f"• {pt}", y_cursor)
            # Constrain bullet circle/text to left column only — increase spacing to avoid overlap
            y_cursor += 0.90

        # Right column — pie chart
        if pie_data:
            cd = ChartData()
            cd.categories = [label for label, _ in pie_data]
            cd.add_series('', [val for _, val in pie_data])

            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                Inches(4.8), Inches(1.20),
                Inches(5.0), Inches(5.50),
                cd
            )
            chart = chart_shape.chart

            # Style the chart
            chart.has_legend = True
            from pptx.enum.chart import XL_LEGEND_POSITION
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            plot = chart.plots[0]
            # Guard for versions of python-pptx where data_labels may not exist
            try:
                if hasattr(plot, 'has_data_labels'):
                    plot.has_data_labels = True
                if hasattr(plot, 'data_labels'):
                    dls = plot.data_labels
                    try:
                        dls.show_percentage = True
                    except Exception:
                        pass
                    try:
                        dls.show_category_name = False
                    except Exception:
                        pass
                    try:
                        dls.show_value = False
                    except Exception:
                        pass
                    try:
                        dls.font.size = Pt(11)
                        dls.font.bold = True
                        dls.font.color.rgb = self.colors["white"]
                    except Exception:
                        pass
            except Exception:
                # Best-effort: if data labels are unsupported, continue without raising
                pass

            # Apply theme colors to slices
            slice_colors = [
                self.colors["primary"],
                self.colors["accent"],
                self.colors["primary_light"],
                RGBColor(16, 185, 129),   # teal
                RGBColor(239, 68, 68),    # red
            ]
            for si, point in enumerate(plot.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = slice_colors[si % len(slice_colors)]

        if page_num:
            self._add_footer(slide, page_num)

    def add_bar_chart_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Bar chart slide — description at top, clustered bar chart below.
        Chart data extracted from numeric key_points.
        """
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.oxml.ns import qn

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Description
        y_cursor = 1.40
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            # Use a smaller chars_per_line estimate so the textbox is taller and avoids clipping
            chars_per_line = 60
            est_lines = max(1, -(-len(short_desc) // chars_per_line))
            ph = max(0.50, est_lines * 0.30 + 0.15)
            pb = slide.shapes.add_textbox(Inches(0.8), Inches(y_cursor), Inches(8.4), Inches(ph))
            ptf = pb.text_frame
            try:
                ptf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(13)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_cursor += ph + 0.20

        # Extract bar data
        bar_data = extract_bar_data(key_points, description)

        if bar_data:
            cd = ChartData()
            cd.categories = [label for label, _ in bar_data]
            cd.add_series('Values', [val for _, val in bar_data])

            chart_top = y_cursor
            chart_h   = 7.0 - chart_top - 0.45   # leave room for footer
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.6), Inches(chart_top),
                Inches(8.8), Inches(chart_h),
                cd
            )
            chart = chart_shape.chart
            chart.has_legend = False

            # Value axis styling — light gridlines for better readability
            val_axis = chart.value_axis
            val_axis.has_major_gridlines = True
            val_axis.major_gridlines.format.line.color.rgb = self.colors["border"]
            try:
                # Format large numbers on axis
                val_axis.tick_labels.font.size = Pt(10)
                val_axis.tick_labels.font.color.rgb = self.colors["text_secondary"]
            except Exception:
                pass

            # Category axis font and styling
            cat_axis = chart.category_axis
            try:
                cat_axis.tick_labels.font.size = Pt(11)
                cat_axis.tick_labels.font.bold = True
                cat_axis.tick_labels.font.color.rgb = self.colors["text_primary"]
            except Exception:
                pass

            # Data labels with better formatting
            plot = chart.plots[0]
            try:
                if hasattr(plot, 'has_data_labels'):
                    plot.has_data_labels = True
                if hasattr(plot, 'data_labels'):
                    dls = plot.data_labels
                    try:
                        dls.show_value = True
                        dls.show_legend_key = False
                    except Exception:
                        pass
                    try:
                        dls.font.size = Pt(10)
                        dls.font.bold = True
                        dls.font.color.rgb = self.colors["white"]
                    except Exception:
                        pass
            except Exception:
                pass

            # Color all bars with primary color and subtle shading
            for idx, point in enumerate(plot.series[0].points):
                point.format.fill.solid()
                # Alternate shades for visual interest
                if idx % 2 == 0:
                    point.format.fill.fore_color.rgb = self.colors["primary"]
                else:
                    point.format.fill.fore_color.rgb = self.colors["primary_light"]

        if page_num:
            self._add_footer(slide, page_num)

    def add_line_chart_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Add line chart slide for trend visualization.
        Shows data progression over time with formatted axis labels.
        """
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Description
        y_cursor = 1.40
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            pb = slide.shapes.add_textbox(Inches(0.8), Inches(y_cursor), Inches(8.4), Inches(0.55))
            ptf = pb.text_frame
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(13)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_cursor += 0.75

        # Extract time series data
        time_labels, time_values = extract_time_series_data(key_points, description)

        if len(time_labels) >= 3:
            cd = ChartData()
            cd.categories = time_labels
            cd.add_series('Trend', time_values)

            chart_top = y_cursor
            chart_h = 7.0 - chart_top - 0.45
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE,
                Inches(0.6), Inches(chart_top),
                Inches(8.8), Inches(chart_h),
                cd
            )
            chart = chart_shape.chart
            chart.has_legend = False

            # Styling
            plot = chart.plots[0]
            try:
                if hasattr(plot, 'has_data_labels'):
                    plot.has_data_labels = True
            except Exception:
                pass

            # Color line
            try:
                plot.series[0].format.line.color.rgb = self.colors["primary"]
                plot.series[0].format.line.width = Pt(2.5)
            except Exception:
                pass

        if page_num:
            self._add_footer(slide, page_num)

    def add_area_chart_slide(self, title: str, key_points: List[str], description: str = "", page_num: int = 0):
        """
        Add stacked area chart for composition or accumulation visualization.
        """
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"]

        # Header accent line
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]

        # Description
        y_cursor = 1.40
        if description:
            sentences = [s.strip() for s in description.replace('!', '.').replace('?', '.').split('.') if s.strip()]
            short_desc = '. '.join(sentences[:2])
            if short_desc and not short_desc.endswith('.'):
                short_desc += '.'
            pb = slide.shapes.add_textbox(Inches(0.8), Inches(y_cursor), Inches(8.4), Inches(0.55))
            ptf = pb.text_frame
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = short_desc
            pp.font.size = Pt(13)
            pp.font.color.rgb = self._body_color()
            pp.line_spacing = 1.3
            y_cursor += 0.75

        # Try to extract multi-series data
        series_data = extract_comparison_series(key_points, description)
        time_labels, time_values = extract_time_series_data(key_points, description)

        if series_data:
            # Multi-series area chart
            cd = ChartData()
            all_categories = []
            for series_values in series_data.values():
                all_categories.extend([label for label, _ in series_values])
            all_categories = sorted(list(set(all_categories)))
            cd.categories = all_categories

            for series_name, series_values in list(series_data.items())[:4]:
                values = []
                for cat in all_categories:
                    val = next((v for l, v in series_values if l == cat), 0)
                    values.append(val)
                cd.add_series(series_name, values)

            chart_top = y_cursor
            chart_h = 7.0 - chart_top - 0.45
            try:
                chart_shape = slide.shapes.add_chart(
                    XL_CHART_TYPE.AREA_STACKED,
                    Inches(0.6), Inches(chart_top),
                    Inches(8.8), Inches(chart_h),
                    cd
                )
                chart = chart_shape.chart

                # Apply theme colors
                slice_colors = [
                    self.colors["primary"],
                    self.colors["accent"],
                    self.colors["primary_light"],
                    RGBColor(16, 185, 129),
                ]
                for si, point in enumerate(chart.plots[0].series):
                    try:
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = slice_colors[si % len(slice_colors)]
                    except Exception:
                        pass
            except Exception:
                pass
        elif len(time_labels) >= 3:
            # Fallback to simple area chart
            cd = ChartData()
            cd.categories = time_labels
            cd.add_series('Data', time_values)

            chart_top = y_cursor
            chart_h = 7.0 - chart_top - 0.45
            try:
                chart_shape = slide.shapes.add_chart(
                    XL_CHART_TYPE.AREA_STACKED,
                    Inches(0.6), Inches(chart_top),
                    Inches(8.8), Inches(chart_h),
                    cd
                )
            except Exception:
                pass

        if page_num:
            self._add_footer(slide, page_num)

    def add_agenda_slide(self, topics: List[str]):
        """
        Add an agenda/overview slide after the title slide.
        Topics are displayed as numbered items in a visually appealing layout.
        """
        # Handle empty topics list
        if not topics or len(topics) == 0:
            return
        
        # Clean up topics (remove empty strings)
        topics = [t.strip() for t in topics if t and str(t).strip()]
        if len(topics) == 0:
            return
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors["white"] if self.theme != "dark" else self.colors["bg_light"]
        
        # Header accent line
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.08))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors["accent"]
        header_shape.line.color.rgb = self.colors["accent"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.35))
        sub_frame = sub_box.text_frame
        sub_p = sub_frame.paragraphs[0]
        sub_p.text = "Here's what we'll cover"
        sub_p.font.size = Pt(14)
        sub_p.font.color.rgb = self.colors["text_secondary"]
        sub_p.font.italic = True
        
        # Determine layout based on number of topics
        num_topics = min(len(topics), 8)
        
        if num_topics <= 4:
            # Single column layout — large agenda items
            self._add_agenda_single_column(slide, topics[:num_topics])
        elif num_topics <= 6:
            # Two-column layout
            self._add_agenda_two_column(slide, topics[:num_topics])
        else:
            # Two-column compact layout
            self._add_agenda_two_column_compact(slide, topics[:num_topics])

    
    def _add_agenda_single_column(self, slide, topics: List[str]):
        """Render agenda items in single column (large format)."""
        y_start = 1.8
        item_height = 1.0
        
        for idx, topic in enumerate(topics):
            y = y_start + idx * item_height
            
            # Numbered circle background
            circle_x = 0.6
            circle_y = y + 0.15
            circle = slide.shapes.add_shape(9, Inches(circle_x), Inches(circle_y), Inches(0.4), Inches(0.4))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors["primary"] if idx % 2 == 0 else self.colors["accent"]
            circle.line.fill.background()
            
            # Number inside circle
            ctf = circle.text_frame
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.text = str(idx + 1)
            cp.font.size = Pt(18)
            cp.font.bold = True
            cp.font.color.rgb = self.colors["white"]
            cp.alignment = PP_ALIGN.CENTER
            
            # Topic text
            text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8.3), Inches(0.7))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_p = text_frame.paragraphs[0]
            text_p.text = topic
            text_p.font.size = Pt(18)
            text_p.font.bold = True
            text_p.font.color.rgb = self.colors["primary"]
            text_p.line_spacing = 1.2
    
    def _add_agenda_two_column(self, slide, topics: List[str]):
        """Render agenda items in two columns."""
        mid = (len(topics) + 1) // 2
        left_topics = topics[:mid]
        right_topics = topics[mid:]
        
        col_positions = [0.6, 5.3]
        y_start = 1.8
        item_height = 0.9
        
        for col_idx, col_topics in enumerate([left_topics, right_topics]):
            col_x = col_positions[col_idx]
            
            for idx, topic in enumerate(col_topics):
                item_idx = idx + (mid if col_idx == 1 else 0)
                y = y_start + idx * item_height
                
                # Numbered circle
                circle = slide.shapes.add_shape(9, Inches(col_x), Inches(y + 0.10), Inches(0.35), Inches(0.35))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors["primary"] if item_idx % 2 == 0 else self.colors["accent"]
                circle.line.fill.background()
                
                ctf = circle.text_frame
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
                cp = ctf.paragraphs[0]
                cp.text = str(item_idx + 1)
                cp.font.size = Pt(14)
                cp.font.bold = True
                cp.font.color.rgb = self.colors["white"]
                cp.alignment = PP_ALIGN.CENTER
                
                # Topic text
                text_box = slide.shapes.add_textbox(Inches(col_x + 0.50), Inches(y), Inches(4.2), Inches(0.65))
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_p = text_frame.paragraphs[0]
                text_p.text = topic
                text_p.font.size = Pt(15)
                text_p.font.bold = True
                text_p.font.color.rgb = self.colors["primary"]
                text_p.line_spacing = 1.1
    
    def _add_agenda_two_column_compact(self, slide, topics: List[str]):
        """Render agenda items in two columns, compact format."""
        mid = (len(topics) + 1) // 2
        left_topics = topics[:mid]
        right_topics = topics[mid:]
        
        col_positions = [0.6, 5.3]
        y_start = 1.8
        item_height = 0.75
        
        for col_idx, col_topics in enumerate([left_topics, right_topics]):
            col_x = col_positions[col_idx]
            
            for idx, topic in enumerate(col_topics):
                item_idx = idx + (mid if col_idx == 1 else 0)
                y = y_start + idx * item_height
                
                # Smaller numbered circle
                circle = slide.shapes.add_shape(9, Inches(col_x), Inches(y + 0.08), Inches(0.30), Inches(0.30))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors["primary"] if item_idx % 2 == 0 else self.colors["accent"]
                circle.line.fill.background()
                
                ctf = circle.text_frame
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
                cp = ctf.paragraphs[0]
                cp.text = str(item_idx + 1)
                cp.font.size = Pt(11)
                cp.font.bold = True
                cp.font.color.rgb = self.colors["white"]
                cp.alignment = PP_ALIGN.CENTER
                
                # Topic text — smaller
                text_box = slide.shapes.add_textbox(Inches(col_x + 0.45), Inches(y), Inches(4.3), Inches(0.55))
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_p = text_frame.paragraphs[0]
                text_p.text = topic
                text_p.font.size = Pt(13)
                text_p.font.bold = True
                text_p.font.color.rgb = self.colors["primary"]
                text_p.line_spacing = 1.1

    def save(self, filename: str = "presentation.pptx"):

        print("SAVE FUNCTION EXECUTED")

        Path("presentations").mkdir(exist_ok=True)

        filepath = f"presentations/{filename}"

        self.prs.save(filepath)

        print("PPT SAVED LOCALLY:", filepath)

        return filepath

# Example usage
if __name__ == "__main__":
    ppt = PresentationGenerator("SlideSmith", "Generate professional presentations in seconds")
    
    ppt.add_title_slide()
    
    ppt.add_content_slide(
        "About SlideSmith",
        [
            "• AI-powered presentation generation",
            "• Powered by Llama3 and LangGraph",
            "• Professional theme-based designs",
            "• PDF to PPTX conversion support"
        ]
    )
    
    ppt.add_comparison_slide(
        "Traditional vs AI Presentations",
        [
            {"title": "Traditional", "content": "Manual Design\nTime Consuming\nLimited Templates"},
            {"title": "SlideSmith", "content": "AI Generated\n2 Minutes\nCustomizable Themes"},
            {"title": "Results", "content": "Professional\nConsistent\nScalable"}
        ]
    )
    
    ppt.add_summary_slide([
        "Fast: Generate presentations in minutes",
        "Professional: Enterprise-grade design",
        "Intelligent: AI-powered content generation"
    ])
    
    filepath = ppt.save("slidesmith_demo.pptx")
    print(f"✓ Presentation saved: {filepath}")
